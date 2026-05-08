# -*- coding: utf-8 -*-
"""BabyNote 개발 정리 PPT 생성기 (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = r"E:\code\BabyNote\docs\babynote_summary.pptx"

# 코랄핑크 톤 팔레트
BG_LIGHT = RGBColor(0xFF, 0xF5, 0xF0)   # 부드러운 코랄크림
BG_DARK = RGBColor(0xA4, 0x3F, 0x45)    # 짙은 코랄 (커버용)
TITLE_FE7D81 = RGBColor(0xFE, 0x7D, 0x81)
STROKE_A43F45 = RGBColor(0xA4, 0x3F, 0x45)
ACCENT_FFB5A7 = RGBColor(0xFF, 0xB5, 0xA7)
MINT_B6E3C9 = RGBColor(0xB6, 0xE3, 0xC9)
TEXT_DARK = RGBColor(0x4A, 0x45, 0x44)
TEXT_MUTED = RGBColor(0x70, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
KOR_FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, font=KOR_FONT,
             italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT_DARK,
                bullet_color=STROKE_A43F45, font=KOR_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet dot
        b = p.add_run()
        b.text = "● "
        b.font.name = font
        b.font.size = Pt(size)
        b.font.color.rgb = bullet_color
        # body
        r = p.add_run()
        r.text = item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def add_pill(slide, x, y, w, h, fill=ACCENT_FFB5A7, line=STROKE_A43F45,
             line_w=1.2):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_stripe(slide):
    """좌측 코랄 핑크 가는 띠 — 모든 본문 슬라이드 공통 액센트."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0), Inches(0.18),
                               Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT_FFB5A7
    s.line.fill.background()
    return s


def page_header(slide, num, total, title, subtitle=None):
    add_stripe(slide)
    add_text(slide, 0.55, 0.32, 9, 0.85, title,
             size=32, bold=True, color=STROKE_A43F45)
    if subtitle:
        add_text(slide, 0.55, 1.12, 9, 0.5, subtitle,
                 size=14, color=TEXT_MUTED, italic=True)
    add_text(slide, 11.4, 0.42, 1.7, 0.4,
             f"{num:02d} / {total:02d}", size=12,
             color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def section_card(slide, x, y, w, h, *, fill=WHITE, line=ACCENT_FFB5A7,
                 line_w=1.2):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


# ===== Slide 1 — 표지 =====
s = prs.slides.add_slide(BLANK)
set_bg(s, BG_LIGHT)
# 큰 코랄 라운드 카드
section_card(s, 1.2, 1.4, 10.9, 4.9, fill=WHITE, line=ACCENT_FFB5A7,
             line_w=2.5)
add_text(s, 1.5, 1.9, 10.3, 1.3, "BabyNote",
         size=72, bold=True, color=TITLE_FE7D81, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.15, 10.3, 0.7, "개발 정리",
         size=32, bold=True, color=STROKE_A43F45, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 4.05, 10.3, 0.5,
         "글로벌 베이비 케어 기록 앱 — KR / JP / EN",
         size=18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
# 버전 pill
add_pill(s, 5.3, 4.95, 2.7, 0.6, fill=ACCENT_FFB5A7, line=STROKE_A43F45,
         line_w=1.5)
add_text(s, 5.3, 4.95, 2.7, 0.6, "v1.0.0",
         size=22, bold=True, color=STROKE_A43F45, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 6.55, 10.3, 0.4, "Flutter · Supabase · Material 3",
         size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)

# 헬퍼: 본문 슬라이드 만들기
def body(num, title, subtitle, bullets, footer=None):
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, BG_LIGHT)
    page_header(sl, num, 17, title, subtitle)
    section_card(sl, 0.55, 1.7, 12.2, 5.4)
    add_bullets(sl, 0.95, 1.95, 11.4, 5.0, bullets, size=15)
    if footer:
        add_text(sl, 0.55, 7.05, 12.2, 0.35, footer,
                 size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
                 italic=True)
    return sl


# ===== Slide 2 — 앱 개요 =====
body(2, "앱 개요", "Newborn~24M 양육 기록 + 분유·기저귀 재고 + 의료 일정",
     [
         "Flutter 3.41+ 단일 코드베이스로 Android 우선, iOS 호환 가능",
         "Supabase Auth + Postgres RLS — 가족 공유 RLS 적용",
         "한국어 / 일본어 / 영어 — ARB 기반 다국어",
         "타깃: 초보 부모, 새벽 한 손 사용 환경 — 큰 버튼 / 큰 글씨",
         "오프라인 우선 X (MVP는 온라인 위주), Phase별 기능 점진 확장",
     ])

# ===== Slide 3 — 홈 화면 =====
body(3, "홈 화면", "한 화면에 핵심 정보 모두 노출 (스크롤 최소화)",
     [
         "AppBar — Baby Note 스트로크 타이틀 (fill #FE7D81 + stroke #A43F45)",
         "자녀 카드 — 이름·생후 일수 + 체중·키·머리둘레 + WHO 백분위 배지(상위 N% / 하위 N%)",
         "오늘의 요약 차트 — 수유 횟수·수면 시간·기저귀 횟수 가로 bar",
         "오늘의 기록 4 타일 — 수유 / 수면 / 기저귀 / 성장 (마지막 활동 시간 + 알림 도트)",
         "데이터/관리 그리드 — 재고 관리 / 기록 편집 / 성장 통계",
         "의료 그리드 — 단골 병원 / 예방접종",
         "간편수유 FAB — 짧은 탭 1회 저장 / 길게 누름 ml 직접 입력",
         "다가오는 일정 종 — 분유·기저귀·접종 알림 모음 (시트 70% 높이 + 스크롤)",
         "백버튼 — 첫 번째 SnackBar 안내, 2초 이내 두 번째에 종료",
     ])

# ===== Slide 4 — 인트로 / 로딩 =====
body(4, "인트로 & 로딩 애니메이션", "Lottie + 자체 CustomPainter 듀얼 트랙",
     [
         "5초 인트로 화면 — 핑크 코끼리 + Baby Note 타이틀 → 자동 진입",
         "ElephantLoader (CustomPainter) — 몸/머리/귀/트렁크/다리 전부 코드 그리기",
         "달리기 모션 — 좌우 sin 이동 + 다리 4개 교차 + 어깨 라인 하이라이트",
         "트렁크 들어올리기 — 사이클 후반 t=0.5~1.0에 위로 컬업",
         "코끝에서 하트 3개 순차 방출(spawn 0.55/0.65/0.75) → 위로 솟으며 페이드아웃",
         "Lottie loading.json 있으면 우선 사용 → 없으면 ElephantLoader fallback",
         "30+ 페이지의 CircularProgressIndicator를 일괄 BabyLoading 으로 교체",
     ])

# ===== Slide 5 — 코치 마크 온보딩 =====
body(5, "코치 마크 온보딩", "최초 진입 시 핵심 6개 영역을 dimmed 배경 + 안내",
     [
         "tutorial_coach_mark 패키지 + GlobalKey 6개로 위젯 위치 추적",
         "안내 순서 — 자녀 추가 → 알림 종 → 오늘의 기록 → 데이터/관리 → 의료 → 간편수유 FAB",
         "포커스 영역은 paddingFocus 0 + 텍스트는 위쪽 ContentAlign.top + extraOffset",
         "스킵 버튼 텍스트: \"다시 보지 않기\" — 우상단 배치(FAB와 겹침 방지)",
         "세션 플래그 _dismissedThisSession — 한 번 닫히면 같은 세션에서 자동 표시 X",
         "설정 → 홈 화면 도움말 다시 보기 → 다음 앱 재실행 시 다시 표시",
     ])

# ===== Slide 6 — 분유 재고 =====
body(6, "분유 재고", "가루분유와 액상분유 듀얼 모드 + 한국 시장 브랜드 프리셋",
     [
         "form 컬럼 (powder | liquid) + 마이그레이션 시 기존은 액상으로 전환",
         "가루분유: 통 무게(g) + 1스쿱 g + 1스쿱 ml — 잔량 계산식 포함",
         "액상분유: 1팩/병 ml + 팩 수량 — 단순 ml 차감",
         "한국 8개 브랜드 프리셋 — 남양 임페리얼XO, 매일 앱솔루트/페어리, 일동 산양, 네슬레 NAN, 압타밀, 힙(HiPP), 거버",
         "브랜드 칩 탭 시 1스쿱 g/ml 자동 채움 (수동 수정 가능)",
         "잔량 표시 — 남은 g·ml + 예상 소진일(D-N) + 색상 경고",
     ])

# ===== Slide 7 — 성장 통계 =====
body(7, "성장 통계 — WHO 표준 비교", "체중·키·머리둘레 백분위 + 5선 곡선 차트",
     [
         "WHO 0~24개월 LMS 파라미터 인라인 번들 (boys/girls × weight/height/HC)",
         "Box-Cox Z-score 공식 + erf/inverse-CDF 근사로 정확한 percentile 계산",
         "사용자 친화 라벨 — \"또래 평균\" / \"상위 N%\" / \"하위 N%\" (P{n} 대신)",
         "color band — 정상 녹색 / 주의 주황 / 외곽 빨강",
         "성장 차트 페이지 — fl_chart로 P3·P15·P50·P85·P97 5선 + 자녀 측정점",
         "가상 아이 크기 시각화 — baby emoji 사이즈를 키(cm)에 비례하게 스케일",
         "의료 진단 아님 면책 표시 + 길게 누름으로 진입",
     ])

# ===== Slide 8 — 기간 통계 =====
body(8, "기간 통계", "1주일 / 1개월 토글 — 수유·수면·기저귀 추세 + 세분화",
     [
         "1주일 모드 — 7일 일별 막대 차트, X축 M/d",
         "1개월 모드 — 12개월 월별 막대 차트, X축 M월 (5일 간격 라벨)",
         "기간 통계 카드 3개 — 코랄/민트/세컨더리 컬러 테두리",
         "수유 — 총 회수/ml + 모유·분유·이유식 세분화",
         "수면 — 총 시간 + 평균/일 또는 /월",
         "기저귀 — 총 회수 + 소변·대변·둘 다 세분화",
         "stats_providers fetch limit 1000 — 1개월 데이터 충분 커버",
     ])

# ===== Slide 9 — 종합 기록 =====
body(9, "종합 기록 (기록 편집)", "수유·수면·기저귀를 시간순 통합한 일별 타임라인",
     [
         "단일 페이지 — 탭 없이 데일리 통합 화면",
         "수유 / 수면 / 기저귀 3종 데이터를 한 리스트에 합쳐서 표시",
         "최신 → 과거 정렬 + yyyy-MM-dd 그룹 헤더 (오늘 / 어제 / 날짜)",
         "각 카드 — 이모지 + 요약 + HH:mm",
         "탭 → 편집 페이지, 길게 누름 → 삭제 confirm",
         "성장 측정 기록은 별도 통계 화면(성장 통계 탭)으로 이동",
     ])

# ===== Slide 10 — 단골 병원 =====
body(10, "단골 병원", "Google Places Autocomplete + 자동 채움",
     [
         "PlacesService (HTTP) — autocomplete + details endpoint 직접 호출",
         "language / country 분기 — KR=kr, JP=jp, EN=null (글로벌)",
         "350ms debounce + 후보 5개 드롭다운",
         "후보 선택 시 — 이름/주소/전화/위경도 자동 채움",
         "API 키는 GOOGLE_PLACES_API_KEY (dart-define) — run/dev.json",
         "키 비어있으면 자동완성 비활성화 (앱 정상 동작)",
     ])

# ===== Slide 11 — 예방접종 =====
body(11, "예방접종", "한국 표준 schedule 기반 + 다가오는 14일 알림",
     [
         "vaccine_schedules — 국가별 표준(KR 우선) seed 데이터",
         "vaccinations — 자녀별 완료 기록 (vaccine_code + dose_number 키)",
         "카드 상태 — 완료 / 다가옴(D-N) / 지난 미접종 / 예정",
         "지난 미접종은 빨간 경고 배지로 강조",
         "홈의 다가오는 일정 종에서 14일 이내 미접종 알림 표시",
         "예방접종 일정 화면에서 한 번에 \"접종 완료\" 기록 가능",
     ])

# ===== Slide 12 — Android 홈 위젯 =====
body(12, "Android 홈 위젯", "4×1 컴팩트 — 잠금화면에서 1탭 진입",
     [
         "수유 / 수면 / 기저귀 / 성장 4 타일 — 이모지 + 라벨 + 요약 한 줄",
         "딥링크 — babynote://widget/{type} → MainActivity → GoRouter redirect",
         "Flutter 측 HomeWidgetPublisher — Riverpod watch + saveWidgetData",
         "데이터 변경 시 자동으로 위젯 SharedPreferences 갱신 + updateWidget 호출",
         "수면 진행 중일 때 라벨 \"수면중\"으로 동적 변경",
         "minHeight 70dp + targetCellHeight 1 — 컴팩트 레이아웃",
     ])

# ===== Slide 13 — 알림 =====
body(13, "알림 시스템", "flutter_local_notifications + ZonedSchedule",
     [
         "분유 잔량 < 3일 — 코랄 errorContainer 강조",
         "기저귀 사이즈업 14일 이내",
         "예방접종 14일 이내 미접종",
         "성장 측정 — 마지막 + 7일 후 09:00 주1회 리마인드",
         "수면 진행 중 ongoing notification (앱 종료해도 유지)",
         "Sentry 통합 — 알림 스케줄러 실패 자동 로깅",
     ])

# ===== Slide 14 — 다국어 + 폰트 =====
body(14, "다국어 + 폰트", "ARB 기반 i18n + Jua(주아체) 둥근 폰트 전면 적용",
     [
         "lib/l10n/app_ko.arb / app_en.arb / app_ja.arb — 모든 라벨 분리 관리",
         "GoogleFonts.juaTextTheme — textTheme 전체에 Jua 적용",
         "AppBar titleTextStyle도 GoogleFonts.jua(...) 직접 적용 (textTheme 자동 안 되므로)",
         "TabBar / SegmentedButton 라벨도 Jua",
         "Korean 라운드 디스플레이 폰트 — 영문 글리프도 둥근 톤",
     ])

# ===== Slide 15 — 디자인 시스템 =====
sl = prs.slides.add_slide(BLANK)
set_bg(sl, BG_LIGHT)
page_header(sl, 15, 17, "디자인 시스템", "파스텔 코랄핑크 + 파스텔 민트 듀얼 시드")
# 컬러 칩들
chips = [
    ("#FFB5A7", "Primary (Coral Pink)", ACCENT_FFB5A7, WHITE),
    ("#B6E3C9", "Tertiary (Mint)", MINT_B6E3C9, TEXT_DARK),
    ("#FFF5F0", "Scaffold BG", BG_LIGHT, TEXT_DARK),
    ("#FE7D81", "Title Fill", TITLE_FE7D81, WHITE),
    ("#A43F45", "Stroke / Dark Coral", STROKE_A43F45, WHITE),
    ("#D06A5C", "Press Feedback", RGBColor(0xD0, 0x6A, 0x5C), WHITE),
]
section_card(sl, 0.55, 1.7, 12.2, 2.4)
for i, (hex_, label, fill, txt) in enumerate(chips):
    col, row = i % 3, i // 3
    cx = 0.85 + col * 4.0
    cy = 1.95 + row * 1.0
    add_pill(sl, cx, cy, 3.6, 0.85, fill=fill, line=STROKE_A43F45, line_w=1)
    add_text(sl, cx, cy + 0.05, 3.6, 0.4, hex_,
             size=15, bold=True, color=txt, align=PP_ALIGN.CENTER)
    add_text(sl, cx, cy + 0.43, 3.6, 0.38, label,
             size=11, color=txt, align=PP_ALIGN.CENTER)
# 추가 설명
section_card(sl, 0.55, 4.2, 12.2, 2.9)
add_bullets(sl, 0.95, 4.45, 11.4, 2.6, [
    "ColorScheme.fromSeed dual seed pattern — primary=코랄, tertiary=민트",
    "라이트 테마 한정 scaffold/canvas/AppBar 배경 통일 (Coral Cream)",
    "모든 Card 기본 — 흰 배경 + 코랄핑크 45% 1px 테두리로 배경과 구분",
    "TabBar 선택=코랄핑크 / 미선택=흐림(45%) / 인디케이터=primary",
    "SegmentedButton 선택=핑크 35% 파스텔 + 다크 코랄 글자 — 모유 좌/우, 가루/액상 등",
    "FilledButton/OutlinedButton — 홈 메뉴 카드와 동일 톤 (흰+코랄 외곽선+다크 코랄 글자)",
    "splashColor / highlightColor — Coral Deep #D06A5C 28% / 16%",
], size=14)

# ===== Slide 16 — 기술 스택 =====
sl = prs.slides.add_slide(BLANK)
set_bg(sl, BG_LIGHT)
page_header(sl, 16, 17, "기술 스택", "주요 패키지 / 인프라")
# 2열 카드
section_card(sl, 0.55, 1.7, 6.0, 5.4)
section_card(sl, 6.85, 1.7, 5.9, 5.4)
add_text(sl, 0.85, 1.85, 5.6, 0.5, "Flutter 패키지",
         size=18, bold=True, color=STROKE_A43F45)
add_bullets(sl, 0.85, 2.4, 5.6, 4.6, [
    "flutter_riverpod 2.x — 상태 관리",
    "go_router — 선언형 라우팅 + redirect",
    "fl_chart — 막대/라인 차트",
    "lottie + 자체 CustomPainter — 로딩",
    "google_fonts (Jua)",
    "tutorial_coach_mark — 코치 마크",
    "home_widget — Android 위젯 sync",
    "flutter_local_notifications + timezone",
    "image_picker · share_plus · url_launcher",
    "supabase_flutter — Auth + Postgres",
    "http — Google Places API 호출",
], size=13)
add_text(sl, 7.15, 1.85, 5.5, 0.5, "백엔드 / 인프라 / 도구",
         size=18, bold=True, color=STROKE_A43F45)
add_bullets(sl, 7.15, 2.4, 5.5, 4.6, [
    "Supabase — Auth, Postgres, RLS",
    "13개 마이그레이션 파일 (init~caregiver_invites)",
    "Sentry — 실시간 에러 모니터링",
    "Google Places API (KR/JP/EN)",
    "WHO 0~24M 성장 표준 LMS 인라인",
    "shared_preferences — 테마/온보딩 플래그",
    "GitHub Actions (옵션) — 향후 CI/CD",
    "Android 무선 디버깅으로 듀얼 폰 동시 설치",
], size=13)

# ===== Slide 17 — 마무리 =====
sl = prs.slides.add_slide(BLANK)
set_bg(sl, BG_LIGHT)
section_card(sl, 1.2, 1.4, 10.9, 4.9, fill=WHITE,
             line=ACCENT_FFB5A7, line_w=2.5)
add_text(sl, 1.5, 2.0, 10.3, 0.85, "Wrap up",
         size=44, bold=True, color=TITLE_FE7D81, align=PP_ALIGN.CENTER)
add_text(sl, 1.5, 2.95, 10.3, 0.55,
         "BabyNote v1.0.0 — 새벽에도 한 손으로",
         size=20, bold=True, color=STROKE_A43F45,
         align=PP_ALIGN.CENTER)
# 향후 계획
add_text(sl, 1.5, 3.7, 10.3, 0.45, "향후 계획",
         size=16, bold=True, color=STROKE_A43F45,
         align=PP_ALIGN.CENTER)
fut = [
    "iOS 위젯 (WidgetKit) + 라이브 액티비티",
    "Push 알림 (FCM) — 가족 공유 상황 동기화",
    "결제 (RevenueCat) — 다자녀 / 클라우드 백업",
    "가족 공유 확장 — 초대 코드 흐름 / 권한 그래뉴얼",
    "오프라인 우선 (drift + Brick) 동기화",
]
add_bullets(sl, 3.4, 4.2, 6.6, 1.7, fut, size=13)
add_text(sl, 1.5, 6.55, 10.3, 0.4, "감사합니다",
         size=16, color=TEXT_MUTED, italic=True,
         align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Saved: {OUT}")
